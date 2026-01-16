"""
Script to load a PPTX presentation into security_presentation collection.

Usage:
    python scripts/load_security_presentation.py <pptx_file_path> [presentation_number] [presentation_name]
"""

import asyncio
import sys
import json
from pathlib import Path
from pptx import Presentation
from motor.motor_asyncio import AsyncIOMotorClient
from pymongo.errors import ConnectionFailure, ServerSelectionTimeoutError, DuplicateKeyError

# Add parent directory to path to import app modules
backend_dir = Path(__file__).parent.parent
sys.path.insert(0, str(backend_dir))

from app.core.config import settings
from app.utils.datetime_utils import utc_now


def extract_pptx_content(file_path: str) -> dict:
    """
    Extract content from a PPTX file.
    
    Args:
        file_path: Path to the PPTX file
        
    Returns:
        Dictionary with extracted content
    """
    print(f"Reading PPTX file: {file_path}")
    
    try:
        prs = Presentation(file_path)
        
        # Extract slides
        slides = []
        full_text_parts = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_content = {
                "slide_number": slide_idx,
                "title": "",
                "shapes": [],
                "text": []
            }
            
            # Extract text from all shapes on the slide
            slide_text_parts = []
            for shape in slide.shapes:
                shape_data = {
                    "shape_type": shape.shape_type,
                    "text": ""
                }
                
                # Check if shape has text
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        shape_data["text"] = text
                        slide_content["shapes"].append(shape_data)
                        slide_text_parts.append(text)
                        
                        # Check if it's a title (usually first text box or placeholder)
                        if not slide_content["title"] and shape.shape_type == 1:  # 1 = PLACEHOLDER
                            slide_content["title"] = text
                
                # Handle tables
                if shape.shape_type == 19:  # 19 = TABLE
                    table_data = {
                        "table_number": len([s for s in slide.shapes if s.shape_type == 19 and hasattr(s, 'table')]) + 1,
                        "rows": []
                    }
                    
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                cell_text = cell.text.strip() if cell.text else ""
                                row_data.append(cell_text)
                            if any(row_data):  # Only add non-empty rows
                                table_data["rows"].append(row_data)
                        
                        if table_data["rows"]:
                            shape_data["table"] = table_data
                            slide_content["shapes"].append(shape_data)
            
            # Combine all text from slide
            slide_text = "\n".join(slide_text_parts)
            if slide_text:
                slide_content["text"] = slide_text
                full_text_parts.append(f"Slide {slide_idx}: {slide_text}")
            
            if slide_content["shapes"] or slide_content["text"]:
                slides.append(slide_content)
        
        # Combine all text
        full_text = "\n\n".join(full_text_parts)
        
        # Extract metadata
        metadata = {
            "file_name": Path(file_path).name,
            "slide_count": len(prs.slides),
            "total_shapes": sum(len(slide["shapes"]) for slide in slides)
        }
        
        return {
            "file_name": Path(file_path).name,
            "slides": slides,
            "slide_count": len(slides),
            "full_text": full_text,
            "total_characters": len(full_text),
            "metadata": metadata
        }
        
    except Exception as e:
        print(f"Error extracting PPTX content: {e}")
        import traceback
        traceback.print_exc()
        raise


async def load_presentation_to_mongodb(
    file_path: str,
    presentation_number: int,
    presentation_name: str
):
    """
    Load a PPTX presentation into MongoDB security_presentation collection.
    
    Args:
        file_path: Path to the PPTX file
        presentation_number: Presentation number (must be unique)
        presentation_name: Presentation name
    """
    # Get connection details from environment variables
    connection_string = os.getenv("DATABASE_URL")
    if not connection_string:
        print("[ERROR] DATABASE_URL environment variable is not set")
        return False
    database_name = os.getenv("DATABASE_NAME", "bsg_demo")
    
    try:
        print(f"Connecting to MongoDB...")
        print(f"Database: {database_name}")
        
        client = AsyncIOMotorClient(
            connection_string,
            serverSelectionTimeoutMS=30000,
            connectTimeoutMS=30000,
            socketTimeoutMS=30000,
        )
        
        # Test connection
        await client.admin.command('ping')
        print("[OK] Connection successful!")
        
        # Get database and collection
        db = client[database_name]
        collection = db["security_presentation"]
        
        # Extract content from PPTX
        print(f"\nExtracting content from PPTX file...")
        presentation_content = extract_pptx_content(file_path)
        
        print(f"[OK] Extracted {presentation_content['slide_count']} slides")
        print(f"[OK] Total characters: {presentation_content['total_characters']}")
        
        # Prepare document for MongoDB
        document = {
            "presentation_number": presentation_number,
            "presentation_name": presentation_name,
            "presentation": presentation_content,
            "created_at": utc_now(),
            "updated_at": utc_now()
        }
        
        # Check if document with this presentation_number already exists
        existing = await collection.find_one({"presentation_number": presentation_number})
        if existing:
            print(f"\n[WARN] Presentation with number {presentation_number} already exists")
            print(f"Updating existing document...")
            result = await collection.update_one(
                {"presentation_number": presentation_number},
                {"$set": {
                    "presentation_name": presentation_name,
                    "presentation": presentation_content,
                    "updated_at": utc_now()
                }}
            )
            print(f"[OK] Updated document: {result.modified_count} document(s) modified")
        else:
            print(f"\nInserting new document...")
            result = await collection.insert_one(document)
            print(f"[OK] Inserted document with ID: {result.inserted_id}")
        
        # Verify insertion
        verify_doc = await collection.find_one({"presentation_number": presentation_number})
        if verify_doc:
            print(f"\n[SUCCESS] Document verified in database:")
            print(f"  - Presentation Number: {verify_doc['presentation_number']}")
            print(f"  - Presentation Name: {verify_doc['presentation_name']}")
            print(f"  - Slides: {verify_doc['presentation'].get('slide_count', 0)}")
            print(f"  - Total Characters: {verify_doc['presentation'].get('total_characters', 0)}")
        
        client.close()
        return True
        
    except DuplicateKeyError as e:
        print(f"[ERROR] Duplicate key error: {e}")
        return False
    except (ConnectionFailure, ServerSelectionTimeoutError) as e:
        print(f"[ERROR] Connection failed: {e}")
        return False
    except Exception as e:
        print(f"[ERROR] Error: {e}")
        import traceback
        traceback.print_exc()
        return False


if __name__ == "__main__":
    import os
    
    if len(sys.argv) < 2:
        print("Usage: python load_security_presentation.py <pptx_file_path> [presentation_number] [presentation_name]")
        sys.exit(1)
    
    file_path = sys.argv[1]
    presentation_number = int(sys.argv[2]) if len(sys.argv) > 2 else 1
    presentation_name = sys.argv[3] if len(sys.argv) > 3 else "Security Presentation"
    
    # Check if file exists
    if not os.path.exists(file_path):
        print(f"[ERROR] File not found: {file_path}")
        sys.exit(1)
    
    success = asyncio.run(load_presentation_to_mongodb(file_path, presentation_number, presentation_name))
    sys.exit(0 if success else 1)

