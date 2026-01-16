"""
HTML5 conversion functions for security presentations.
"""

from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from app.core.logging import get_logger

logger = get_logger(__name__)


def convert_pptx_to_html5(pptx_path: str) -> str:
    """
    Convert PPTX file to HTML5 format.
    
    Args:
        pptx_path: Path to PPTX file
        
    Returns:
        HTML5 string representation of the presentation
    """
    try:
        prs = Presentation(pptx_path)
        html_parts = []
        
        # Start HTML5 document
        html_parts.append('<!DOCTYPE html>')
        html_parts.append('<html lang="en">')
        html_parts.append('<head>')
        html_parts.append('<meta charset="UTF-8">')
        html_parts.append('<meta name="viewport" content="width=device-width, initial-scale=1.0">')
        html_parts.append('<title>Presentation</title>')
        html_parts.append('<style>')
        html_parts.append('''
            * {
                margin: 0;
                padding: 0;
                box-sizing: border-box;
            }
            body {
                font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif;
                line-height: 1.6;
                color: #333;
                background-color: #f5f5f5;
                padding: 20px;
            }
            .presentation-container {
                max-width: 1200px;
                margin: 0 auto;
                background: white;
                border-radius: 8px;
                box-shadow: 0 2px 8px rgba(0,0,0,0.1);
                padding: 30px;
            }
            .slide {
                background: white;
                margin: 30px 0;
                padding: 40px;
                border-radius: 8px;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                min-height: 400px;
                page-break-after: always;
            }
            .slide-number {
                font-size: 14px;
                color: #666;
                margin-bottom: 15px;
                font-weight: 600;
                text-transform: uppercase;
                letter-spacing: 1px;
            }
            .slide-title {
                font-size: 28px;
                font-weight: bold;
                color: #283054;
                margin-bottom: 25px;
                border-bottom: 3px solid #283054;
                padding-bottom: 15px;
            }
            .slide-content {
                font-size: 16px;
                line-height: 1.8;
                color: #4A5568;
            }
            .slide-content p {
                margin: 15px 0;
            }
            .slide-content ul, .slide-content ol {
                margin: 15px 0;
                padding-left: 30px;
            }
            .slide-content li {
                margin: 8px 0;
            }
            table {
                width: 100%;
                border-collapse: collapse;
                margin: 25px 0;
                background: white;
                box-shadow: 0 1px 3px rgba(0,0,0,0.1);
            }
            table th, table td {
                border: 1px solid #ddd;
                padding: 15px;
                text-align: left;
            }
            table th {
                background-color: #283054;
                color: white;
                font-weight: 600;
                font-size: 14px;
            }
            table tr:nth-child(even) {
                background-color: #f9f9f9;
            }
            table tr:hover {
                background-color: #f0f0f0;
            }
            .table-title {
                font-weight: 600;
                margin-top: 25px;
                margin-bottom: 12px;
                color: #283054;
                font-size: 18px;
            }
            .shape-content {
                margin: 15px 0;
                padding: 10px;
                background: #f8f9fa;
                border-left: 4px solid #283054;
            }
            .metadata {
                background: #e8f4f8;
                padding: 20px;
                border-radius: 5px;
                margin-top: 40px;
                font-size: 14px;
            }
            .metadata h3 {
                margin-top: 0;
                color: #283054;
                margin-bottom: 15px;
            }
            .metadata p {
                margin: 8px 0;
            }
        ''')
        html_parts.append('</style>')
        html_parts.append('</head>')
        html_parts.append('<body>')
        html_parts.append('<div class="presentation-container">')
        
        # Process slides
        for slide_idx, slide in enumerate(prs.slides, 1):
            html_parts.append('<div class="slide">')
            
            # Slide number
            html_parts.append(f'<div class="slide-number">Slide {slide_idx}</div>')
            
            # Extract title (usually first placeholder or text box)
            slide_title = ""
            slide_content_parts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        # Check if it's a title placeholder
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            if not slide_title:
                                slide_title = text
                                continue
                        
                        # Add text content
                        if text not in slide_content_parts:
                            slide_content_parts.append(text)
                
                # Handle tables
                if shape.has_table:
                    table_html = ['<div class="table-title">Table</div>', '<table>']
                    for row_idx, row in enumerate(shape.table.rows):
                        table_html.append('<tr>')
                        for cell in row.cells:
                            cell_tag = 'th' if row_idx == 0 else 'td'
                            cell_text = cell.text.strip() if cell.text else ""
                            table_html.append(f'<{cell_tag}>{cell_text}</{cell_tag}>')
                        table_html.append('</tr>')
                    table_html.append('</table>')
                    slide_content_parts.append(''.join(table_html))
            
            # Add slide title
            if slide_title:
                html_parts.append(f'<div class="slide-title">{slide_title}</div>')
            
            # Add slide content
            html_parts.append('<div class="slide-content">')
            for content in slide_content_parts:
                if content.startswith('<table>'):
                    html_parts.append(content)
                else:
                    # Split into paragraphs
                    paragraphs = content.split('\n\n')
                    for para in paragraphs:
                        if para.strip():
                            # Check if it's a list
                            if para.strip().startswith('-') or para.strip().startswith('•'):
                                lines = para.strip().split('\n')
                                html_parts.append('<ul>')
                                for line in lines:
                                    if line.strip().startswith('-') or line.strip().startswith('•'):
                                        item = line.strip().lstrip('-•').strip()
                                        if item:
                                            html_parts.append(f'<li>{item}</li>')
                                html_parts.append('</ul>')
                            else:
                                html_parts.append(f'<p>{para.strip().replace(chr(10), "<br>")}</p>')
            html_parts.append('</div>')
            
            html_parts.append('</div>')
        
        # Add metadata
        html_parts.append('<div class="metadata">')
        html_parts.append('<h3>Presentation Metadata</h3>')
        html_parts.append(f'<p><strong>Total Slides:</strong> {len(prs.slides)}</p>')
        html_parts.append(f'<p><strong>File:</strong> {Path(pptx_path).name}</p>')
        html_parts.append('</div>')
        
        # Close HTML
        html_parts.append('</div>')
        html_parts.append('</body>')
        html_parts.append('</html>')
        
        return '\n'.join(html_parts)
        
    except ImportError as e:
        raise Exception(f"Required library not installed: {e}")
    except Exception as e:
        logger.error(f"Error converting PPTX to HTML5: {e}", exc_info=True)
        raise

